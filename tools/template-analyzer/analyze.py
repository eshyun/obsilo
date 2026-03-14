"""
Template Analyzer -- Core Pipeline

Analyzes a PPTX template and generates a Visual Design Language (VDL) Skill document.

Pipeline:
1. python-pptx: Extract shapes, positions, text, theme
2. LibreOffice headless: Render slides to PNG
3. Spatial analysis: Detect composition patterns
4. Claude Vision: Multimodal interpretation
5. Output: SKILL.md in VDL format
"""

from __future__ import annotations

import base64
import io
import json
import os
import re
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import openai
from PIL import Image
from pptx import Presentation
from pptx.util import Emu


# ---------------------------------------------------------------------------
#  Data classes
# ---------------------------------------------------------------------------

@dataclass
class ShapeData:
    """Extracted shape information from python-pptx."""
    name: str
    shape_id: int
    left: int  # EMU
    top: int   # EMU
    width: int  # EMU
    height: int  # EMU
    text: str
    placeholder_type: str | None
    placeholder_idx: int | None
    has_text_frame: bool
    geometry: str  # e.g. "rect", "roundRect", "chevron", "custGeom"
    fill_type: str
    is_group: bool


@dataclass
class SlideData:
    """All extracted data for a single slide."""
    number: int  # 1-based
    layout_name: str
    shapes: list[ShapeData]
    image_path: str | None = None  # Path to rendered PNG


@dataclass
class BrandDNA:
    """Theme colors and fonts."""
    colors: dict[str, str]  # name -> hex
    fonts: dict[str, str]   # major, minor
    slide_width: int = 12192000   # EMU
    slide_height: int = 6858000   # EMU


@dataclass
class CompositionAnalysis:
    """Claude Vision analysis result for a composition."""
    classification: str
    composition_name: str
    meaning: str
    use_when: str
    visual_quality: str
    slide_numbers: list[int]
    shapes: list[dict[str, Any]]


@dataclass
class AnalysisResult:
    """Complete analysis result."""
    template_name: str
    slide_count: int
    brand_dna: BrandDNA
    slides: list[SlideData]
    compositions: list[CompositionAnalysis] = field(default_factory=list)


# ---------------------------------------------------------------------------
#  Step 1: python-pptx extraction
# ---------------------------------------------------------------------------

def extract_structural_data(pptx_path: str) -> tuple[list[SlideData], BrandDNA]:
    """Extract shapes, theme, and layout info from PPTX."""
    prs = Presentation(pptx_path)

    # Brand DNA from theme
    brand_dna = _extract_brand_dna(prs)

    slides: list[SlideData] = []
    for i, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        shapes = _extract_shapes(slide)
        slides.append(SlideData(number=i, layout_name=layout_name, shapes=shapes))

    return slides, brand_dna


def _extract_brand_dna(prs: Presentation) -> BrandDNA:
    """Extract colors and fonts from the theme."""
    colors: dict[str, str] = {}
    fonts = {"major": "Calibri", "minor": "Calibri"}

    slide_width = prs.slide_width or Emu(12192000)
    slide_height = prs.slide_height or Emu(6858000)

    # Try to extract theme colors from the XML
    try:
        theme = prs.slide_masters[0].element
        # Navigate to theme XML
        from lxml import etree
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        # Theme is in the slide master's related parts
        for rel in prs.slide_masters[0].part.rels.values():
            if "theme" in rel.reltype:
                theme_xml = rel.target_part.blob
                tree = etree.fromstring(theme_xml)

                # Colors
                color_names = [
                    "dk1", "dk2", "lt1", "lt2",
                    "accent1", "accent2", "accent3", "accent4",
                    "accent5", "accent6", "hlink", "folHlink",
                ]
                for cname in color_names:
                    elems = tree.xpath(f".//a:clrScheme/a:{cname}//a:srgbClr", namespaces=ns)
                    if elems:
                        colors[cname] = f"#{elems[0].get('val', '000000')}"
                    else:
                        # Try system color
                        sys_elems = tree.xpath(f".//a:clrScheme/a:{cname}//a:sysClr", namespaces=ns)
                        if sys_elems:
                            colors[cname] = f"#{sys_elems[0].get('lastClr', '000000')}"

                # Fonts
                major_elems = tree.xpath(".//a:majorFont/a:latin", namespaces=ns)
                if major_elems:
                    fonts["major"] = major_elems[0].get("typeface", "Calibri")
                minor_elems = tree.xpath(".//a:minorFont/a:latin", namespaces=ns)
                if minor_elems:
                    fonts["minor"] = minor_elems[0].get("typeface", "Calibri")
                break
    except Exception:
        pass  # Fall back to defaults

    return BrandDNA(
        colors=colors,
        fonts=fonts,
        slide_width=int(slide_width),
        slide_height=int(slide_height),
    )


def _extract_shapes(slide) -> list[ShapeData]:
    """Extract all shapes from a slide."""
    shapes: list[ShapeData] = []
    for shape in slide.shapes:
        geom = "rect"
        try:
            if hasattr(shape, "auto_shape_type"):
                ast = shape.auto_shape_type
                if ast is not None:
                    geom = str(ast).split("(")[0].strip()
        except (ValueError, AttributeError):
            pass
        if geom == "rect" and shape.shape_type is not None:
            st = str(shape.shape_type)
            if "FREEFORM" in st or "CUSTOM" in st.upper():
                geom = "custGeom"
            elif "GROUP" in st:
                geom = "group"
            else:
                geom = st.split("(")[0].strip().lower()

        text = ""
        has_text_frame = False
        if shape.has_text_frame:
            has_text_frame = True
            text = shape.text_frame.text

        fill_type = "none"
        try:
            if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
                fill_type = str(shape.fill.type).split("(")[0].strip().lower()
        except Exception:
            pass

        ph_type = None
        ph_idx = None
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type = str(ph.type).split("(")[0].strip() if ph.type is not None else None
            ph_idx = ph.idx

        shapes.append(ShapeData(
            name=shape.name,
            shape_id=shape.shape_id,
            left=int(shape.left or 0),
            top=int(shape.top or 0),
            width=int(shape.width or 0),
            height=int(shape.height or 0),
            text=text[:200],  # Truncate
            placeholder_type=ph_type,
            placeholder_idx=ph_idx,
            has_text_frame=has_text_frame,
            geometry=geom,
            fill_type=fill_type,
            is_group=shape.shape_type is not None and "GROUP" in str(shape.shape_type),
        ))

    return shapes


# ---------------------------------------------------------------------------
#  Step 2: LibreOffice rendering
# ---------------------------------------------------------------------------

def render_slides_to_png(pptx_path: str, output_dir: str) -> list[str]:
    """Render each slide to PNG using LibreOffice headless."""
    # Find LibreOffice
    lo_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        "/snap/bin/libreoffice",
    ]
    soffice = None
    for p in lo_paths:
        if os.path.exists(p):
            soffice = p
            break
    if soffice is None:
        # Try PATH
        try:
            result = subprocess.run(["which", "soffice"], capture_output=True, text=True)
            if result.returncode == 0:
                soffice = result.stdout.strip()
        except Exception:
            pass

    if soffice is None:
        raise RuntimeError(
            "LibreOffice not found. Install it:\n"
            "  macOS: brew install --cask libreoffice\n"
            "  Linux: sudo apt install libreoffice\n"
        )

    # Convert to PNG
    subprocess.run(
        [soffice, "--headless", "--convert-to", "png", "--outdir", output_dir, pptx_path],
        capture_output=True,
        timeout=120,
    )

    # LibreOffice only exports the first slide as single PNG.
    # For multi-slide, we convert to PDF first, then PDF pages to PNG.
    pdf_dir = os.path.join(output_dir, "pdf_tmp")
    os.makedirs(pdf_dir, exist_ok=True)

    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", pdf_dir, pptx_path],
        capture_output=True,
        timeout=120,
    )

    # Find the PDF
    pdf_files = list(Path(pdf_dir).glob("*.pdf"))
    if not pdf_files:
        raise RuntimeError("LibreOffice PDF conversion failed")

    pdf_path = str(pdf_files[0])

    # Convert PDF pages to PNG using Pillow (requires poppler or pdf2image)
    png_paths = _pdf_to_pngs(pdf_path, output_dir)

    return png_paths


def _pdf_to_pngs(pdf_path: str, output_dir: str) -> list[str]:
    """Convert PDF pages to individual PNGs."""
    png_paths: list[str] = []

    # Try PyMuPDF first (no external binary needed)
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(dpi=150)
            out = os.path.join(output_dir, f"slide_{i:03d}.png")
            pix.save(out)
            png_paths.append(out)
        doc.close()
        return png_paths
    except ImportError:
        pass

    # Fallback: pdf2image (needs poppler binaries)
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path, dpi=150, fmt="png")
        for i, img in enumerate(images, 1):
            out = os.path.join(output_dir, f"slide_{i:03d}.png")
            img.save(out, "PNG")
            png_paths.append(out)
        return png_paths
    except (ImportError, Exception):
        pass

    raise RuntimeError(
        "Cannot convert PDF to PNG. Install one of:\n"
        "  pip install PyMuPDF\n"
        "  pip install pdf2image  (+ poppler: brew install poppler)\n"
    )


# ---------------------------------------------------------------------------
#  Step 3: Spatial analysis (heuristic pre-classification)
# ---------------------------------------------------------------------------

def classify_slide_heuristic(slide: SlideData) -> str:
    """Heuristic slide classification based on shape positions and types."""
    shapes = slide.shapes
    text_shapes = [s for s in shapes if s.has_text_frame and s.text.strip()]
    placeholders = [s for s in shapes if s.placeholder_type]

    if len(shapes) == 0 or (len(text_shapes) == 0 and len(placeholders) == 0):
        return "blank"

    # Title slide: few shapes, has title placeholder, centered
    title_phs = [s for s in placeholders if s.placeholder_type and "TITLE" in s.placeholder_type.upper()]
    subtitle_phs = [s for s in placeholders if s.placeholder_type and "SUBTITLE" in s.placeholder_type.upper()]
    if title_phs and subtitle_phs and len(text_shapes) <= 4:
        return "title"

    # Section: large centered text, dark background
    if title_phs and not subtitle_phs and len(text_shapes) <= 2:
        return "section"

    # Chevrons / process shapes
    chevron_shapes = [s for s in shapes if "chevron" in s.geometry.lower() or "homePlate" in s.geometry.lower()]
    if len(chevron_shapes) >= 3:
        return "process"

    # KPI: multiple small, evenly-spaced shapes with short text
    small_text_shapes = [s for s in text_shapes if len(s.text) < 30 and s.width < 4000000]
    if len(small_text_shapes) >= 4:
        # Check if they're in a grid pattern
        tops = sorted(set(s.top for s in small_text_shapes))
        if len(tops) <= 3:
            return "kpi"

    # Two-column: shapes split left/right
    if len(text_shapes) >= 2:
        slide_mid = 6096000  # Half of standard width
        left_shapes = [s for s in text_shapes if s.left + s.width / 2 < slide_mid]
        right_shapes = [s for s in text_shapes if s.left + s.width / 2 >= slide_mid]
        if left_shapes and right_shapes and abs(len(left_shapes) - len(right_shapes)) <= 2:
            return "two-column"

    # Pyramid: triangular shape or stacked decreasing-width shapes
    pyramid_shapes = [s for s in shapes if "triangle" in s.geometry.lower() or "pyramid" in s.name.lower()]
    if pyramid_shapes:
        return "pyramid"

    # Default: content
    return "content"


# ---------------------------------------------------------------------------
#  Step 4: Claude Vision analysis
# ---------------------------------------------------------------------------

def analyze_with_vision(
    client: openai.OpenAI,
    slides: list[SlideData],
    brand_dna: BrandDNA,
    vision_prompt: str,
    model: str = "anthropic/claude-sonnet-4-5",
    on_progress: Any = None,
) -> list[CompositionAnalysis]:
    """Send slide images + structural data to Claude Vision via OpenRouter."""
    batch_size = 10
    all_analyses: list[dict] = []

    total_batches = (len(slides) + batch_size - 1) // batch_size

    for batch_idx in range(0, len(slides), batch_size):
        batch = slides[batch_idx:batch_idx + batch_size]
        current_batch = batch_idx // batch_size + 1

        if on_progress:
            on_progress(f"Analysiere Batch {current_batch}/{total_batches} (Slides {batch[0].number}-{batch[-1].number})...")

        # Build OpenAI-format content parts
        content_parts: list[dict] = []

        # Structural context
        structural_data = []
        for slide in batch:
            slide_info = {
                "slide_number": slide.number,
                "layout_name": slide.layout_name,
                "heuristic_classification": classify_slide_heuristic(slide),
                "shapes": [
                    {
                        "name": s.name,
                        "geometry": s.geometry,
                        "left": s.left,
                        "top": s.top,
                        "width": s.width,
                        "height": s.height,
                        "text": s.text[:100],
                        "placeholder_type": s.placeholder_type,
                        "has_text": s.has_text_frame,
                    }
                    for s in slide.shapes
                ],
            }
            structural_data.append(slide_info)

        content_parts.append({
            "type": "text",
            "text": f"## Structural Data\n\n```json\n{json.dumps(structural_data, indent=2)}\n```",
        })

        # Slide images (OpenAI vision format: image_url with data URI)
        for slide in batch:
            if slide.image_path and os.path.exists(slide.image_path):
                # Resize to max 800px wide to reduce token usage
                img = Image.open(slide.image_path)
                if img.width > 800:
                    ratio = 800 / img.width
                    img = img.resize((800, int(img.height * ratio)), Image.LANCZOS)
                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=75)
                img_data = buf.getvalue()
                b64 = base64.standard_b64encode(img_data).decode("utf-8")
                content_parts.append({
                    "type": "text",
                    "text": f"\n## Slide {slide.number} ({slide.layout_name})",
                })
                content_parts.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{b64}",
                    },
                })

        content_parts.append({
            "type": "text",
            "text": "\nAnalyze these slides and return the JSON array as specified.",
        })

        try:
            response = client.chat.completions.create(
                model=model,
                max_tokens=16384,
                messages=[
                    {"role": "system", "content": vision_prompt},
                    {"role": "user", "content": content_parts},
                ],
            )
        except Exception as api_err:
            if on_progress:
                on_progress(f"API-Fehler bei Batch {current_batch}: {api_err}")
            continue

        # Parse JSON from response
        if not response.choices:
            # OpenRouter may return error in a different field
            error_msg = getattr(response, "error", None) or "Keine Antwort vom Modell"
            if on_progress:
                on_progress(f"Warnung Batch {current_batch}: {error_msg}")
            continue

        choice = response.choices[0]
        response_text = choice.message.content or ""
        was_truncated = getattr(choice, "finish_reason", None) == "length"

        if was_truncated and on_progress:
            on_progress(f"Warnung Batch {current_batch}: Antwort wurde abgeschnitten (max_tokens erreicht)")

        batch_analyses = _extract_json_array(response_text)
        if batch_analyses is not None:
            all_analyses.extend(batch_analyses)
            if on_progress:
                on_progress(f"Batch {current_batch}: {len(batch_analyses)} Slides analysiert")
        else:
            if on_progress:
                # Show first 200 chars of response for debugging
                preview = response_text[:200].replace("\n", " ")
                on_progress(f"Warnung Batch {current_batch}: JSON-Parsing fehlgeschlagen. Antwort: {preview}...")

    # Group into compositions
    return _group_into_compositions(all_analyses, slides)


def _extract_json_array(text: str) -> list[dict] | None:
    """Extract a JSON array from LLM response text, handling various formats."""
    text = text.strip()

    # Try 1: Extract from ```json ... ``` code block (complete)
    code_block = re.search(r"```(?:json)?\s*(\[[\s\S]*?\])\s*```", text)
    if code_block:
        try:
            return json.loads(code_block.group(1))
        except json.JSONDecodeError:
            pass

    # Try 2: Find the outermost [ ... ] using bracket-counting
    start = text.find("[")
    if start == -1:
        return None

    depth = 0
    end = -1
    in_string = False
    escape_next = False
    for i in range(start, len(text)):
        c = text[i]
        if escape_next:
            escape_next = False
            continue
        if c == "\\":
            escape_next = True
            continue
        if c == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if c == "[":
            depth += 1
        elif c == "]":
            depth -= 1
            if depth == 0:
                end = i + 1
                break

    if end != -1:
        try:
            return json.loads(text[start:end])
        except json.JSONDecodeError:
            pass

    # Try 3: Incomplete/truncated JSON -- repair by closing open structures
    snippet = text[start:]
    return _repair_truncated_json_array(snippet)


def _repair_truncated_json_array(text: str) -> list[dict] | None:
    """Attempt to repair a truncated JSON array by closing open structures.

    Strategy: Find the last complete object in the array, discard the
    incomplete trailing object, and close the array.
    """
    # Find the last complete }, then close the array
    # We search backwards for a } that could end a complete object
    last_complete = -1
    depth_brace = 0
    depth_bracket = 0
    in_str = False
    esc = False

    for i in range(len(text)):
        c = text[i]
        if esc:
            esc = False
            continue
        if c == "\\":
            esc = True
            continue
        if c == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if c == "{":
            depth_brace += 1
        elif c == "}":
            depth_brace -= 1
            if depth_brace == 0 and depth_bracket <= 1:
                # This } closes a top-level object in the array
                last_complete = i
        elif c == "[":
            depth_bracket += 1
        elif c == "]":
            depth_bracket -= 1

    if last_complete == -1:
        return None

    # Take text up to and including the last complete object, close the array
    repaired = text[: last_complete + 1].rstrip().rstrip(",") + "]"
    try:
        result = json.loads(repaired)
        if isinstance(result, list):
            return result
    except json.JSONDecodeError:
        pass

    return None


def _shape_signature(slide_data: dict) -> str:
    """Create a structural fingerprint from a slide's shape list.

    Two slides with the same classification and similar shape structure
    (same placeholder types and rough count of text shapes) are considered
    the same composition, regardless of what the LLM named them.
    """
    shapes = slide_data.get("shapes", [])
    # Collect placeholder roles and shape roles
    roles = sorted(s.get("role", s.get("shape_name", "?")) for s in shapes)
    return "|".join(roles)


def _group_into_compositions(
    analyses: list[dict],
    all_slides: list[SlideData] | None = None,
) -> list[CompositionAnalysis]:
    """Group individual slide analyses into compositions.

    Grouping strategy (layered):
    1. Same classification (e.g. "chart", "process")
    2. Same layout_name from the PPTX (strongest signal for identical layouts)
    3. Similar shape structure (fingerprint of shape roles)

    This produces far fewer compositions than grouping by LLM-generated names.
    """
    # Build a lookup: slide_number -> layout_name
    layout_by_slide: dict[int, str] = {}
    if all_slides:
        for s in all_slides:
            layout_by_slide[s.number] = s.layout_name

    # Phase 1: Group by classification + layout_name (from PPTX metadata)
    groups: dict[str, list[dict]] = {}
    for a in analyses:
        cls = a.get("classification", "content")
        slide_num = a.get("slide_number", 0)
        layout = layout_by_slide.get(slide_num, "unknown")
        # Primary grouping: classification + layout name
        group_key = f"{cls}::{layout}"
        if group_key not in groups:
            groups[group_key] = []
        groups[group_key].append(a)

    # Phase 2: Within each group, if shape signatures differ significantly,
    # split into sub-groups (handles cases where same layout has variants)
    final_groups: dict[str, list[dict]] = {}
    for group_key, items in groups.items():
        if len(items) <= 1:
            final_groups[group_key] = items
            continue

        # Sub-group by shape signature
        sub_groups: dict[str, list[dict]] = {}
        for item in items:
            sig = _shape_signature(item)
            if sig not in sub_groups:
                sub_groups[sig] = []
            sub_groups[sig].append(item)

        if len(sub_groups) == 1:
            final_groups[group_key] = items
        else:
            for i, (sig, sub_items) in enumerate(sub_groups.items()):
                final_groups[f"{group_key}::v{i}"] = sub_items

    # Build composition objects
    compositions: list[CompositionAnalysis] = []
    for group_key, items in final_groups.items():
        # Pick the representative with the best visual quality
        quality_rank = {"high": 0, "medium": 1, "low": 2}
        items.sort(key=lambda x: quality_rank.get(x.get("visual_quality", "medium"), 1))
        representative = items[0]
        slide_numbers = sorted(a["slide_number"] for a in items)

        compositions.append(CompositionAnalysis(
            classification=representative.get("classification", "content"),
            composition_name=representative.get("composition_name", "Unknown"),
            meaning=representative.get("meaning", ""),
            use_when=representative.get("use_when", ""),
            visual_quality=representative.get("visual_quality", "medium"),
            slide_numbers=slide_numbers,
            shapes=representative.get("shapes", []),
        ))

    # Sort: title first, then section, then by number of slides (desc)
    priority = {"title": 0, "section": 1}
    compositions.sort(key=lambda c: (
        priority.get(c.classification, 2),
        -len(c.slide_numbers),
    ))

    return compositions


# ---------------------------------------------------------------------------
#  Step 5: Generate VDL output
# ---------------------------------------------------------------------------

def _truncate(text: str, max_len: int) -> str:
    """Truncate text to max_len, appending '...' if shortened."""
    if len(text) <= max_len:
        return text
    return text[: max_len - 3].rstrip() + "..."


def _get_placeholder_shapes(
    slide_data: SlideData | None,
    vision_shapes: list[dict[str, Any]],
) -> list[tuple[str, str, int]]:
    """Return (shape_name, role, max_chars) for replaceable shapes.

    Only includes shapes that have a text frame. Deduplicates by name.
    Limits to the most important shapes (placeholders first, then others).
    """
    if not slide_data:
        return []

    seen: set[str] = set()
    result: list[tuple[str, str, int]] = []

    # Build vision lookup
    vision_by_name = {s.get("shape_name", ""): s for s in vision_shapes}

    # Placeholders first (most important), then other text shapes
    shapes_sorted = sorted(
        slide_data.shapes,
        key=lambda s: (0 if s.placeholder_type else 1, s.name),
    )

    for shape in shapes_sorted:
        if not shape.has_text_frame or shape.name in seen:
            continue
        seen.add(shape.name)

        # Determine role
        vs = vision_by_name.get(shape.name)
        if vs:
            role = vs.get("role", shape.placeholder_type or "content")
            max_chars = vs.get("max_chars", 0)
        else:
            role = shape.placeholder_type or "content"
            max_chars = 0

        # Skip decorative/unnamed shapes that Vision couldn't identify
        if role == shape.name and not shape.placeholder_type:
            continue

        result.append((shape.name, role, max_chars))

    return result[:8]  # Max 8 shapes per composition


def generate_vdl(result: AnalysisResult) -> str:
    """Generate Visual Design Language SKILL.md from analysis result.

    Format is optimized for compactness (<16k chars) while preserving
    all information the LLM needs to create presentations:
    - Which slide to clone (template_slide number)
    - Which shapes to fill (content keys)
    - When to use which composition (semantic meaning)
    """
    name_slug = re.sub(r"\s+", "-", result.template_name.lower())
    trigger_slug = re.sub(r"\s+", "|", result.template_name.lower())
    content_comps = [c for c in result.compositions if c.classification != "blank"]

    lines: list[str] = []

    # YAML frontmatter
    lines.append("---")
    lines.append(f"name: {name_slug}")
    lines.append(f"description: {result.template_name} -- {result.slide_count} Slides, {len(content_comps)} Compositions")
    lines.append(f"trigger: {trigger_slug}")
    lines.append("source: user")
    lines.append("requiredTools: [create_pptx]")
    lines.append("---")
    lines.append("")

    lines.append(f"# {result.template_name} -- Visual Design Language")
    lines.append("")

    # Brand DNA (compact)
    dna = result.brand_dna
    primary = dna.colors.get("dk1", dna.colors.get("accent1", "#000000"))
    accents = ", ".join(
        dna.colors[k] for k in ["accent1", "accent2", "accent3"] if k in dna.colors
    )
    lines.append(f"## Brand-DNA")
    lines.append(f"Primary: {primary} | Accent: {accents} | Heading: {dna.fonts.get('major', 'Calibri')} | Body: {dna.fonts.get('minor', 'Calibri')}")
    lines.append("")

    # Design rules (before compositions -- context for reading)
    lines.append("## Rules")
    lines.append("- ALWAYS use `template_file` + `template_slide` + `content` (NEVER `html`)")
    lines.append("- Shape names in `content` must match EXACTLY (case-sensitive)")
    lines.append("- Never repeat same composition on consecutive slides")
    lines.append("- Max 30% text-only slides")
    lines.append("")

    # Compositions -- compact table format
    lines.append("## Compositions")
    lines.append("")
    lines.append("Each composition: name, slide numbers, when to use, and shape mapping.")
    lines.append("Use the FIRST slide number as `template_slide`. All slides in a group are interchangeable.")
    lines.append("")

    for comp in content_comps:
        nums = comp.slide_numbers
        rep_num = nums[0]

        # Compact slide list
        if len(nums) > 5:
            nums_str = f"{', '.join(str(n) for n in nums[:5])}, +{len(nums) - 5} more"
        else:
            nums_str = ", ".join(str(n) for n in nums)

        # Header with slide numbers
        lines.append(f"### {comp.composition_name} [{comp.classification}]")
        lines.append(f"Slides: {nums_str}")
        lines.append(f"Use: {_truncate(comp.use_when, 120)}")

        # Shape mapping -- inline compact format
        rep_slide = next((s for s in result.slides if s.number == rep_num), None)
        placeholders = _get_placeholder_shapes(rep_slide, comp.shapes)

        if placeholders:
            shape_parts = []
            for name, role, max_chars in placeholders:
                if max_chars:
                    shape_parts.append(f'"{name}": {role} ({max_chars}ch)')
                else:
                    shape_parts.append(f'"{name}": {role}')
            lines.append(f"Shapes (slide {rep_num}): {' | '.join(shape_parts)}")

        lines.append("")

    # Narrative phase mapping (compact)
    lines.append("## Narrative Phases")

    opening_types = {"title", "kpi"}
    tension_types = {"comparison", "two-column", "matrix"}
    resolution_types = {"process", "pyramid", "timeline"}

    phase_map = [
        ("Opening", opening_types, "Set thesis, establish facts"),
        ("Tension", tension_types, "Build contrast, show gap"),
        ("Resolution", resolution_types, "Show path forward"),
    ]

    for phase_name, types, rationale in phase_map:
        comps = [c.composition_name for c in content_comps if c.classification in types]
        if comps:
            lines.append(f"- **{phase_name}**: {', '.join(comps[:6])} -- {rationale}")

    versatile = [
        c.composition_name for c in content_comps
        if c.classification not in opening_types | tension_types | resolution_types
        and c.classification not in {"blank", "image"}
    ]
    if versatile:
        lines.append(f"- **Any phase**: {', '.join(versatile[:8])}")

    lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
#  Main pipeline
# ---------------------------------------------------------------------------

def run_pipeline(
    pptx_path: str,
    api_key: str,
    model: str = "anthropic/claude-sonnet-4-5",
    on_progress: Any = None,
) -> tuple[str, AnalysisResult]:
    """
    Run the complete analysis pipeline.

    Returns (vdl_content, analysis_result).
    """
    pptx_path = os.path.abspath(pptx_path)
    template_name = Path(pptx_path).stem.replace("_", " ").replace("-", " ").strip()

    if on_progress:
        on_progress("Extrahiere Strukturdaten aus PPTX...")

    # Step 1: Extract structural data
    slides, brand_dna = extract_structural_data(pptx_path)

    if on_progress:
        on_progress(f"{len(slides)} Slides gefunden. Rendere mit LibreOffice...")

    # Step 2: Render slides to PNG
    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            png_paths = render_slides_to_png(pptx_path, tmp_dir)
            for i, png_path in enumerate(png_paths):
                if i < len(slides):
                    slides[i].image_path = png_path
            if on_progress:
                on_progress(f"{len(png_paths)} Slide-Bilder gerendert.")
        except RuntimeError as e:
            if on_progress:
                on_progress(f"Warnung: {e}\nFahre ohne Bilder fort (nur Strukturdaten).")

        # Step 3: Heuristic classification (done inline during vision)

        # Step 4: Claude Vision analysis
        if on_progress:
            on_progress("Starte Claude Vision Analyse...")

        # Load vision prompt
        prompt_path = Path(__file__).parent / "prompts" / "vision-prompt.md"
        vision_prompt = prompt_path.read_text(encoding="utf-8")

        client = openai.OpenAI(
            base_url="https://openrouter.ai/api/v1",
            api_key=api_key,
            default_headers={
                "HTTP-Referer": "https://obsilo.ai",
                "X-Title": "Obsilo Template Analyzer",
            },
        )

        has_images = any(s.image_path for s in slides)
        compositions: list[CompositionAnalysis] = []

        if has_images:
            compositions = analyze_with_vision(
                client, slides, brand_dna, vision_prompt, model, on_progress,
            )

        # Fallback: if Vision returned nothing (all batches failed or no images)
        if not compositions:
            if on_progress:
                reason = "Vision-Analyse lieferte keine Ergebnisse" if has_images else "Keine Bilder verfuegbar"
                on_progress(f"{reason} -- nutze heuristische Klassifikation.")
            compositions = _heuristic_only_compositions(slides)

    # Build result
    result = AnalysisResult(
        template_name=template_name,
        slide_count=len(slides),
        brand_dna=brand_dna,
        slides=slides,
        compositions=compositions,
    )

    # Step 5: Generate VDL
    if on_progress:
        on_progress("Generiere Visual Design Language Document...")

    vdl = generate_vdl(result)

    if on_progress:
        on_progress(f"Fertig! VDL: {len(vdl)} Zeichen, {len(compositions)} Kompositionen.")

    return vdl, result


def _heuristic_only_compositions(slides: list[SlideData]) -> list[CompositionAnalysis]:
    """Build compositions from heuristic classification alone (no Vision)."""
    from collections import defaultdict
    groups: dict[str, list[SlideData]] = defaultdict(list)

    for slide in slides:
        classification = classify_slide_heuristic(slide)
        groups[classification].append(slide)

    meaning_map = {
        "title": "Establishes the presentation topic",
        "section": "Marks a new chapter or topic transition",
        "content": "Presents detailed information or bullet points",
        "kpi": "Highlights key metrics at a glance",
        "process": "Shows sequential steps in a workflow",
        "comparison": "Contrasts two or more options",
        "two-column": "Presents two related content areas side by side",
        "table": "Displays structured data in rows and columns",
        "chart": "Visualizes data trends or distributions",
        "pyramid": "Shows hierarchical layers from broad to specific",
        "matrix": "Maps items along two dimensions",
        "timeline": "Shows events or milestones in chronological order",
        "image": "Features visual content prominently",
        "blank": "Empty placeholder slide",
    }

    use_when_map = {
        "title": "Opening or closing a presentation",
        "section": "Transitioning between major topics",
        "content": "Explaining concepts with text and bullets",
        "kpi": "Presenting 2-6 key metrics or numbers",
        "process": "Describing a workflow or sequential steps",
        "comparison": "Contrasting options, pros/cons, before/after",
        "two-column": "Presenting two parallel concepts",
        "table": "Showing structured multi-dimensional data",
        "chart": "Visualizing trends, distributions, or comparisons",
        "pyramid": "Showing hierarchy or priority layers",
        "matrix": "Analyzing items along two axes (e.g., SWOT)",
        "timeline": "Presenting chronological events or milestones",
        "image": "Featuring a key visual or photo",
        "blank": "Reserved for custom content",
    }

    compositions: list[CompositionAnalysis] = []
    for classification, slide_list in groups.items():
        compositions.append(CompositionAnalysis(
            classification=classification,
            composition_name=classification.replace("-", " ").title(),
            meaning=meaning_map.get(classification, ""),
            use_when=use_when_map.get(classification, ""),
            visual_quality="medium",
            slide_numbers=sorted(s.number for s in slide_list),
            shapes=[],
        ))

    priority = {"title": 0, "section": 1}
    compositions.sort(key=lambda c: (priority.get(c.classification, 2), -len(c.slide_numbers)))

    return compositions


# ---------------------------------------------------------------------------
#  CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Analyze a PPTX template and generate a VDL Skill")
    parser.add_argument("template", help="Path to the .pptx template file")
    parser.add_argument("--api-key", required=True, help="Anthropic API key")
    parser.add_argument("--model", default="anthropic/claude-sonnet-4-5", help="OpenRouter model ID")
    parser.add_argument("--output", "-o", default=None, help="Output path for SKILL.md")
    args = parser.parse_args()

    def progress(msg: str) -> None:
        print(f"  {msg}")

    print(f"Analyzing: {args.template}")
    vdl_content, _ = run_pipeline(args.template, args.api_key, args.model, on_progress=progress)

    output_path = args.output or Path(args.template).stem + "-skill.md"
    Path(output_path).write_text(vdl_content, encoding="utf-8")
    print(f"\nSkill written to: {output_path}")
    print(f"Size: {len(vdl_content)} chars")
